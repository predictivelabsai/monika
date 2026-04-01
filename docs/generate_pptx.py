"""Generate AHMF Platform Overview PowerPoint from screenshots."""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).parent.parent
GUIDE = ROOT / "static" / "guide"
OUT = ROOT / "docs" / "AHMF_Platform_Overview.pptx"

# Ashland Hill blue
AH_BLUE = RGBColor(0x00, 0x66, 0xCC)
AH_DARK = RGBColor(0x1E, 0x29, 0x3B)
AH_GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Blue gradient background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = AH_BLUE

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = WHITE
        p2.space_before = Pt(20)

    # AH logo text
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(3), Inches(0.5))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "Ashland Hill Media Finance"
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE
    p3.font.bold = True

    return slide


def add_content_slide(title, bullets, screenshot=None, two_col=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Header bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AH_BLUE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AH_DARK

    if screenshot and not two_col:
        # Text on left, screenshot on right
        text_width = Inches(5.5)
        text_left = Inches(0.8)
        img_left = Inches(6.8)
        img_width = Inches(5.8)
    elif two_col:
        text_width = Inches(11.5)
        text_left = Inches(0.8)
        img_left = img_width = None
    else:
        text_width = Inches(11.5)
        text_left = Inches(0.8)
        img_left = img_width = None

    # Bullets
    txBox2 = slide.shapes.add_textbox(text_left, Inches(1.3), text_width, Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()

        if bullet.startswith("**") and bullet.endswith("**"):
            # Bold subheader
            p.text = bullet.strip("*")
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = AH_DARK
            p.space_before = Pt(12)
        elif bullet.startswith("---"):
            p.text = ""
            p.space_before = Pt(8)
        else:
            p.text = bullet
            p.font.size = Pt(14)
            p.font.color.rgb = AH_GRAY
            p.space_before = Pt(4)
            if bullet.startswith("|"):
                p.font.size = Pt(12)
                p.font.name = "Consolas"

    # Screenshot
    if screenshot and img_left:
        img_path = GUIDE / screenshot
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), img_left, Inches(1.3), img_width)

    # Footer
    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(5), Inches(0.3))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Ashland Hill Media Finance — Private & Confidential"
    p3.font.size = Pt(8)
    p3.font.color.rgb = AH_GRAY

    return slide


# ---- Build Slides ----

# 1. Title
add_title_slide(
    "Technology & AI-Driven Solutions",
    "Products Roadmap & Platform Demo — FY2026 Q1"
)

# 2. Vision
add_content_slide("Platform Vision", [
    "**From Single-Picture Lender to Integrated Media Platform**",
    "---",
    "AHMF is a centralized, AI-driven operating system purpose-built for private credit and film finance.",
    "",
    "Replaces fragmented spreadsheets and siloed workflows with a unified data architecture",
    "9 product modules — all operational",
    "AI-powered intelligence across the full financing lifecycle",
    "Real-time portfolio visibility and deal management",
    "Powered by LangGraph + XAI Grok-3 AI engine",
], screenshot="01_welcome.png")

# 3. Film Financing OS Overview
add_content_slide("Film Financing Operating System", [
    "**The core operating system — 6 fully operational modules**",
    "---",
    "Deals — Pipeline dashboard, deal structuring, performance tracking",
    "Contacts — Counterparty profiles, relationship mapping",
    "Sales & Collections — Territory contracts, MG tracking, collection payments",
    "Credit Rating — AI-powered counterparty scoring (AAA-CCC tiers)",
    "Accounting — Transaction ledger, disbursements, repayments, net position",
    "Communications — Deal-linked messages, tasks with deadlines, completion tracking",
], screenshot="02_deals.png")

# 4. Deal Structuring
add_content_slide("Deal Structuring", [
    "**Comprehensive deal creation with all financing parameters**",
    "---",
    "Title, genre, project type, status (Pipeline/Active/Closed)",
    "Loan amount, interest rate, term in months",
    "Budget, borrower, producer, director, cast",
    "Territory and collateral type",
    "Full audit trail with created_by tracking",
], screenshot="03_deal_new.png")

# 4b. Sales & Collections
add_content_slide("Sales & Collections", [
    "**Territory-based revenue tracking and collateral monitoring**",
    "---",
    "Sales contracts linked to deals and distributors",
    "Territory-based MG (Minimum Guarantee) commitments",
    "Collection payment recording with due dates",
    "Status tracking: pending, received, overdue",
    "Variance analysis: projected vs actual collections",
    "Receivable flagging for delayed payments",
], screenshot="05_sales.png")

# 4c. Credit Rating
add_content_slide("Credit Rating", [
    "**AI-powered counterparty strength assessment**",
    "---",
    "Select any distributor, producer, or sales agent",
    "AI generates comprehensive credit profile:",
    "  Credit Score (0-100)",
    "  Payment Reliability (0-100)",
    "  Risk Tier (AAA through CCC)",
    "",
    "**Factor Analysis:**",
    "  Track Record — historical deal performance",
    "  Financial Stability — company financials",
    "  Market Position — competitive standing",
    "  Payment History — on-time payment patterns",
], screenshot="17_credit.png")

# 4d. Accounting
add_content_slide("Accounting", [
    "**Transaction-level financial integrity**",
    "---",
    "Full transaction ledger across all deals",
    "Transaction types: Disbursement, Repayment, Fee, Interest, Adjustment",
    "Multi-currency support (USD, EUR, GBP, CAD, AUD, JPY, CNY)",
    "Net position tracking (repaid + fees - disbursed)",
    "Counterparty linking for audit trail",
    "Reference notes per transaction",
], screenshot="18_accounting.png")

# 4e. Communications
add_content_slide("Communications", [
    "**Execution control layer for deal teams**",
    "---",
    "Deal-linked messages, tasks, and notifications",
    "Task assignment with due dates and deadlines",
    "Interactive checkboxes to mark task completion",
    "Overdue task flagging for missed deadlines",
    "Status dashboard: open tasks, completed, overdue",
    "Full message history per deal",
], screenshot="19_comms.png")

# 5. Sales Estimates
add_content_slide("Sales Estimates Generator", [
    "**Quantitative valuation intelligence for greenlight and collateral appraisal**",
    "---",
    "Script ingestion with PDF parsing + metadata extraction",
    "Genre + comp benchmarking via TMDB & OMDB databases",
    "Territory-level MG projections",
    "Box office forecasting (domestic + international)",
    "Confidence score and model reliability rating",
    "Comparable film analysis with budget, revenue, and cast data",
], screenshot="06_estimates.png")

# 6. Risk Scoring
add_content_slide("Production Risk Scoring", [
    '**"Moody\'s for execution risk"**',
    "---",
    "AI evaluates 6 risk dimensions on a 0-100 scale:",
    "  Script Complexity — locations, stunts, VFX density",
    "  Budget Feasibility — budget vs comparable films",
    "  Schedule Risk — shoot days vs scope",
    "  Jurisdictional Risk — labor laws, permits, stability",
    "  Crew/Talent Risk — availability, reliability, concentration",
    "  Completion Risk — on-time, on-budget probability",
    "",
    "Output: Overall score, risk tier, mitigation recommendations",
], screenshot="08_risk_form.png")

# 7. Smart Budget
add_content_slide("Smart Budgeting Tool", [
    "**Generative budgeting logic + live cost intelligence**",
    "---",
    "AI generates 3 scenarios: Low / Mid / High",
    "8-12 detailed line items per scenario",
    "",
    "**Budget Categories:**",
    "Above-the-Line (talent, director, producer fees)",
    "Below-the-Line Production (crew, equipment, locations)",
    "Below-the-Line Post (editing, VFX, sound, music)",
    "Insurance & Legal",
    "Financing Costs",
    "Contingency (typically 10%)",
], screenshot="09_budget.png")

# 8. Scheduling
add_content_slide("Automated Production Scheduling", [
    "**Predictive production efficiency software**",
    "---",
    "AI generates day-by-day shooting schedules",
    "Location clustering optimization (minimize company moves)",
    "Day/night scene grouping",
    "Actor availability and labor law compliance",
    "Bottleneck prediction",
    "Scenario comparison (schedule A vs B)",
    "Integration with budgeting tool",
], screenshot="10_schedule.png")

# 9. Soft Funding
add_content_slide("Soft Funding Discovery Engine", [
    '**"Kayak for film incentives"**',
    "---",
    "Database of 16+ global incentive programs",
    "11 countries covered",
    "Average rebate: 27.8%",
    "",
    "**Top Programs:**",
    "  Colombia Cash Rebate — 40%",
    "  Ireland Section 481 — 32%",
    "  Georgia Tax Credit — 30%",
    "  Hungary Cash Rebate — 30%",
    "  Australia PDV Offset — 30%",
    "",
    "Built-in rebate calculator and eligibility filters",
], screenshot="04_funding.png")

# 10. Data Room
add_content_slide("Deal Closing & Data Room", [
    "**Infrastructure layer for streamlined closings**",
    "---",
    "Auto-generated 20-item closing checklists",
    "",
    "**6 Categories:**",
    "  Legal (5 items) — loan agreements, security, completion guarantee",
    "  Insurance (3 items) — E&O, production insurance, completion bond",
    "  Financial (4 items) — chain of title, budget, lab access, CAMA",
    "  Distribution (3 items) — agreements, delivery schedule, MG commitments",
    "  Tax Incentives (2 items) — applications, auditor opinion",
    "  Compliance (3 items) — KYC/AML, OFAC, board approval",
    "",
    "Interactive checkbox tracking with progress bars",
], screenshot="13_checklist.png")

# 11. Audience
add_content_slide("Audience & Marketing Intelligence", [
    "**Campaign simulator meets audience heat map**",
    "---",
    "AI-powered audience segmentation (3-5 segments)",
    "Demographic + psychographic clustering with % share",
    "Marketing channel allocation and spend estimates",
    "Release window optimization",
    "Platform strategy (theatrical, streaming, hybrid)",
    "Festival positioning recommendations",
    "Comparable film marketing analysis",
], screenshot="11_audience.png")

# 12. Talent
add_content_slide("Talent Intelligence", [
    "**Taste graph meets market intelligence**",
    "---",
    "Live TMDB actor/director search with popularity data",
    "AI cast recommendations scored on:",
    "  Heat Index (1-10) — current market popularity",
    "  Genre Fit (1-10) — alignment with project tone",
    "  Salary Tier — Low / Mid / High / Premium",
    "  International Sales Impact — Low / Mid / High",
    "",
    "Package simulation: test cast combinations for projected revenue",
    "Comparable role analysis from filmography data",
], screenshot="14_talent_search.png")

# 13. AI Chat
add_content_slide("AI Chat — The Intelligence Layer", [
    "**18 AI tools accessible via natural language or structured commands**",
    "---",
    "deal:list / portfolio — Deal pipeline and portfolio analytics",
    "contact:search / credit:NAME — Contacts and credit ratings",
    "sales:list / transactions / messages — Financial operations",
    "incentives — Search global film tax incentives",
    "talent:search NAME — Look up actors/directors via TMDB",
    "risk:new / budget:new / schedule:new — AI analysis tools",
    "audience:new — Audience & marketing prediction",
    "",
    "Or ask any question in natural language — AI routes to the right tool",
], screenshot="15_chat_help.png")

# 14. Technical
add_content_slide("Technical Summary", [
    "**Production-grade architecture**",
    "---",
    "Frontend: FastHTML + HTMX (server-rendered, no React/JS frameworks)",
    "AI Engine: LangGraph + XAI Grok-3 (18 structured tools)",
    "Database: PostgreSQL (28 tables in ahmf schema)",
    "Film Data: TMDB + OMDB APIs (budgets, revenue, cast, ratings)",
    "Auth: Email/password + JWT session management",
    "Deployment: Docker + Coolify (auto-deploy on git push)",
    "Test Suite: 30 automated tests covering all modules",
    "",
    "**All 9 roadmap products fully operational — no placeholders**",
], two_col=True)

# 15. Next Steps
add_content_slide("Next Steps", [
    "**Continued platform depth and integrations**",
    "---",
    "Sales Estimate Pipeline: 5-node LangGraph pipeline (analyze > comps > territory > forecast > report)",
    "PDF Script Ingestion: Upload screenplays for automated analysis",
    "Closing Workflow Automation: Drawdown requests, digital signatures",
    "Portfolio Analytics Dashboard: IRR tracking, collection monitoring, variance analysis",
    "Investor Reporting: Automated LP/capital partner reports",
    "Mobile Responsive: Tablet-optimized layout for on-the-go deal review",
], two_col=True)

# 16. Contact
add_title_slide(
    "Thank You",
    "Joe Simpson — Managing Partner\nSimon Williams — Managing Partner\n\nAshland Hill Media Finance"
)

# Save
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
